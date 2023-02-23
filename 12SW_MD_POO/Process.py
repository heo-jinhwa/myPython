import os
from pptx import Presentation

def PPT_RPA_Process(RPAobject):
    pptFile = Presentation(os.path.join(RPAobject._master_path, RPAobject.get_master()[0])) # 마스터파일에서 가지고오기

    pptMatch = {
        'header' : ['헤드라인', '헤드라인 메세지', 'date'],
        'sub-body' : ['PART NO', 'P/NAME', 'Set 수', '재질', '두께', 'D/H', '프레스 톤 수', 'BL/PI', '차종','폭, 피치 여유율']
    }

    slide_list = pptFile.slides
    for slide in slide_list:
        shape_list = slide.shapes
        shape_idx = {}
        
        for idx, v in enumerate(shape_list):
            shape_idx[v.name] = idx
        
        for x in pptMatch['header']:
            name_shape = shape_list[shape_idx[x]]
            name_shape.text = '김진호바보23123'

    pptFile.save("RPA 적용 양식.pptx")